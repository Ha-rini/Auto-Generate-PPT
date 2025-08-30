# MultipleFiles/pptx_generator.py
from pptx import Presentation
from pptx.util import Inches
# Removed 'openai' import as it's no longer directly used here
import os
from .llm_services import generate_slide_content # Import the service

def generate_presentation(content: str, guidance: str, api_key: str, template_path: str) -> str:
    # The api_key parameter is now redundant here as llm_services loads it from env
    # However, we keep it in the signature for consistency with main.py's call,
    # though it won't be used directly in this function.

    # Create a new presentation based on the uploaded template
    prs = Presentation(template_path)

    # Combine content and guidance for the LLM prompt
    full_prompt = f"Generate content for a presentation. Focus on: {content}. Additional guidance: {guidance}. Provide the content in a way that can be split into multiple slides, with each slide's content separated by a clear delimiter (e.g., '---SLIDE_BREAK---')."

    # Call the Aipipe service to generate content
    raw_generated_content = generate_slide_content(full_prompt)

    if raw_generated_content.startswith("Error"):
        # Handle error from LLM service
        print(f"Failed to generate content: {raw_generated_content}")
        # Optionally, raise an exception or return a specific error indicator
        raise Exception(f"Failed to generate presentation content: {raw_generated_content}")

    # Process the response and add slides
    # We'll assume the LLM is instructed to use a specific delimiter
    slides_content = raw_generated_content.split("---SLIDE_BREAK---")

    # Ensure there's at least one slide, even if the split doesn't work as expected
    if not slides_content or all(not s.strip() for s in slides_content):
        slides_content = ["No specific slide content generated. Please check the prompt or LLM response."]

    for i, slide_text in enumerate(slides_content):
        slide_text = slide_text.strip()
        if not slide_text: # Skip empty slides if delimiter creates them
            continue

        # Use a layout that has a title and content placeholder (layout index 1 is common)
        # You might need to inspect your template to find appropriate layout indices
        # For a blank slide, use prs.slide_layouts[6] and add text boxes manually
        try:
            slide_layout = prs.slide_layouts[1] # Title and Content layout
        except IndexError:
            print("Warning: Slide layout 1 not found. Using a blank layout.")
            slide_layout = prs.slide_layouts[6] # Blank layout

        slide = prs.slides.add_slide(slide_layout)

        # Add title and content
        if slide_layout.name == "Title and Content" or len(slide.shapes.placeholders) > 1:
            title = slide.shapes.title
            content_box = slide.shapes.placeholders[1] # Assuming placeholder 1 is for content

            title.text = f"Slide {i+1}" # Dynamic title
            content_box.text = slide_text
        else:
            # If using a blank layout or placeholders are not as expected, add text boxes
            left = top = width = height = Inches(1) # Example dimensions
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = f"Slide {i+1}"
            
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(6))
            content_frame = content_box.text_frame
            content_frame.text = slide_text


    # Save the presentation
    output_path = "generated_presentation.pptx"
    prs.save(output_path)

    return output_path

